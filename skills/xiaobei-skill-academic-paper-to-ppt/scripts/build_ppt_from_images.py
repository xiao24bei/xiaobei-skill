#!/usr/bin/env python3
"""Build a 16:9 image-only PowerPoint deck from slide images."""

from __future__ import annotations

import argparse
import re
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


SUPPORTED_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
DEFAULT_OUTPUT = "defense_presentation.pptx"


def natural_key(path: Path) -> list[object]:
    parts = re.split(r"(\d+)", path.name.lower())
    return [int(part) if part.isdigit() else part for part in parts]


def collect_images(image_dir: Path) -> list[Path]:
    if not image_dir.exists():
        raise FileNotFoundError(f"Image directory does not exist: {image_dir}")
    if not image_dir.is_dir():
        raise NotADirectoryError(f"Image path is not a directory: {image_dir}")

    images = [
        path
        for path in image_dir.iterdir()
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    ]
    if not images:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"No slide images found in {image_dir}. Supported formats: {supported}")

    return sorted(images, key=natural_key)


def normalized_image_path(path: Path, temp_dir: Path) -> Path:
    if path.suffix.lower() != ".webp":
        return path

    output = temp_dir / f"{path.stem}.png"
    with Image.open(path) as image:
        image.convert("RGBA").save(output, "PNG")
    return output


def image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        width, height = image.size
    if width <= 0 or height <= 0:
        raise ValueError(f"Invalid image dimensions for {path}: {width}x{height}")
    return width, height


def add_full_slide_image(slide, image_path: Path, slide_width: int, slide_height: int) -> None:
    width_px, height_px = image_size(image_path)
    slide_ratio = slide_width / slide_height
    image_ratio = width_px / height_px

    if image_ratio >= slide_ratio:
        pic_height = slide_height
        pic_width = int(slide_height * image_ratio)
        left = int((slide_width - pic_width) / 2)
        top = 0
    else:
        pic_width = slide_width
        pic_height = int(slide_width / image_ratio)
        left = 0
        top = int((slide_height - pic_height) / 2)

    slide.shapes.add_picture(
        str(image_path),
        left,
        top,
        width=pic_width,
        height=pic_height,
    )


def build_ppt_from_images(image_dir: Path, output: Path) -> None:
    images = collect_images(image_dir)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    output.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="ppt_slide_images_") as temp_root:
        temp_dir = Path(temp_root)
        for image in images:
            slide = presentation.slides.add_slide(blank_layout)
            normalized = normalized_image_path(image, temp_dir)
            add_full_slide_image(
                slide,
                normalized,
                presentation.slide_width,
                presentation.slide_height,
            )

    presentation.save(output)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create a 16:9 PowerPoint deck from full-slide images."
    )
    parser.add_argument(
        "--image-dir",
        required=True,
        type=Path,
        help="Directory containing slide images in png, jpg, jpeg, or webp format.",
    )
    parser.add_argument(
        "--output",
        default=DEFAULT_OUTPUT,
        type=Path,
        help=f"Output PPTX path. Defaults to {DEFAULT_OUTPUT}.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    try:
        build_ppt_from_images(args.image_dir, args.output)
    except Exception as exc:
        print(f"ERROR: {exc}")
        return 1

    print(f"Created PowerPoint deck: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
